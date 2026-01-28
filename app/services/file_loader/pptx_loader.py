from __future__ import annotations

from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.exceptions import ServiceException
from app.services.file_loader.base import FileLoader, ImageInfo, PageContent


def _extract_table_text(table) -> str:
    # 表格内容转为文本
    rows: List[str] = []
    for row in table.rows:
        values: List[str] = []
        for cell in row.cells:
            cell_text = " ".join(cell.text.split())
            values.append(cell_text)
        if any(values):
            rows.append("\t".join(values))
    return "\n".join(rows)


def _parse_pptx(file_path: Path) -> List[PageContent]:
    presentation = Presentation(str(file_path))
    pages: List[PageContent] = []
    image_index = 1
    picture_types = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }

    for slide_index, slide in enumerate(presentation.slides, start=1):
        # 每一张幻灯片视作一页
        texts: List[str] = []
        images: List[ImageInfo] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)
                if text.strip():
                    texts.append(text.strip())
            if shape.has_table:
                table_text = _extract_table_text(shape.table)
                if table_text.strip():
                    texts.append(table_text.strip())
            if shape.shape_type in picture_types:
                try:
                    image = shape.image
                    images.append(
                        ImageInfo(
                            index=image_index,
                            name=getattr(image, "filename", None),
                            width=getattr(image, "width", None),
                            height=getattr(image, "height", None),
                            mime_type=getattr(image, "content_type", None),
                        )
                    )
                    image_index += 1
                except Exception:
                    continue
        pages.append(PageContent(page_number=slide_index, text="\n".join(texts).strip(), images=images))
    return pages


def _parse_with_unstructured(file_path: Path) -> List[PageContent]:
    try:
        from unstructured.partition.auto import partition
    except Exception as exc:
        raise ServiceException("缺少 unstructured 依赖，无法解析该格式") from exc

    try:
        elements = partition(filename=str(file_path))
    except Exception as exc:
        raise ServiceException(f"解析文件失败: {exc}") from exc

    pages_map: Dict[int, PageContent] = {}
    for element in elements:
        # unstructured 的 page_number 可能缺失，默认归为第 1 页
        page_number = getattr(element.metadata, "page_number", None) or 1
        text = element.text or ""
        page = pages_map.setdefault(page_number, PageContent(page_number=page_number, text=""))
        if text:
            if page.text:
                page.text += "\n"
            page.text += text

    pages = [pages_map[key] for key in sorted(pages_map)]
    return pages


class PptxLoader(FileLoader):
    """PPT 解析器，支持 pptx（优先）及 ppt（降级解析）。"""

    def parse(self, file_path: Path) -> List[PageContent]:
        suffix = file_path.suffix.lower()
        if suffix == ".pptx":
            return _parse_pptx(file_path)
        if suffix == ".ppt":
            return _parse_with_unstructured(file_path)
        raise ServiceException(f"不支持的 PPT 格式: {suffix}")
