from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.core.exception.exceptions import ServiceException
from app.rag.file_loader.parsers.base import BaseParser
from app.rag.file_loader.types import ParsedPage


def _extract_table_text(table) -> str:
    """
    功能描述:
        将 PPT 表格内容转为文本，行之间使用换行分隔，列之间使用制表符分隔。

    参数说明:
        table: python-pptx 的表格对象。

    返回值:
        str: 提取后的表格文本。

    异常说明:
        无。
    """
    rows: list[str] = []
    for row in table.rows:
        values: list[str] = []
        for cell in row.cells:
            values.append(" ".join(cell.text.split()))
        if any(values):
            rows.append("\t".join(values))
    return "\n".join(rows)


def _parse_pptx(file_path: Path) -> list[ParsedPage]:
    """
    功能描述:
        解析 pptx 文件，每一张幻灯片视为一页并提取文本/表格内容。

    参数说明:
        file_path (Path): pptx 文件路径。

    返回值:
        list[ParsedPage]: 按幻灯片组织的页面列表。

    异常说明:
        Exception: 幻灯片读取失败时由 python-pptx 抛出。
    """
    presentation = Presentation(str(file_path))
    pages: list[ParsedPage] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    paragraph.text for paragraph in shape.text_frame.paragraphs
                ).strip()
                if text:
                    texts.append(text)
            if shape.has_table:
                table_text = _extract_table_text(shape.table).strip()
                if table_text:
                    texts.append(table_text)
        pages.append(ParsedPage(page_number=slide_index, text="\n".join(texts)))
    return pages


def _parse_ppt(file_path: Path) -> list[ParsedPage]:
    """
    功能描述:
        使用 unstructured 降级解析 ppt 文件并按页聚合文本。

    参数说明:
        file_path (Path): ppt 文件路径。

    返回值:
        list[ParsedPage]: 按页组织的页面列表。

    异常说明:
        ServiceException:
            - 缺少 unstructured 依赖时抛出；
            - unstructured 解析失败时抛出。
    """
    try:
        from unstructured.partition.auto import partition
    except Exception as exc:
        raise ServiceException("缺少 unstructured 依赖，无法解析 ppt 文件") from exc

    try:
        elements = partition(filename=str(file_path))
    except Exception as exc:
        raise ServiceException(f"解析 ppt 文件失败: {exc}") from exc

    pages_map: dict[int, ParsedPage] = {}
    for element in elements:
        page_number = getattr(element.metadata, "page_number", None) or 1
        text = element.text or ""
        page = pages_map.setdefault(page_number, ParsedPage(page_number=page_number))
        if text:
            if page.text:
                page.text += "\n"
            page.text += text
    return [pages_map[index] for index in sorted(pages_map)]


class PptParser(BaseParser):
    """
    功能描述:
        解析 PPT 文件，支持 pptx 与 ppt。

    参数说明:
        无。解析参数通过 `parse` 方法传入。

    返回值:
        无。调用 `parse` 时返回页面列表。

    异常说明:
        ServiceException: 文件后缀不支持或降级解析失败时抛出。
    """

    def parse(self, file_path: Path) -> list[ParsedPage]:
        """
        功能描述:
            根据后缀分发 PPT 解析分支。

        参数说明:
            file_path (Path): PPT 文件路径。

        返回值:
            list[ParsedPage]: 按页组织的页面列表。

        异常说明:
            ServiceException: 不支持的 PPT 文件格式时抛出。
        """
        suffix = file_path.suffix.lower()
        if suffix == ".pptx":
            return _parse_pptx(file_path)
        if suffix == ".ppt":
            return _parse_ppt(file_path)
        raise ServiceException(f"不支持的 PPT 格式: {suffix}")
