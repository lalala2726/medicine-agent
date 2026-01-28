from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.exceptions import ServiceException
from app.core.file_loader.base import (
    FileLoader,
    ImageInfo,
    PageContent,
    ensure_image_output_dir,
    save_image_bytes,
)


def _extract_table_text(table) -> str:
    """
    将 PPT 表格内容转为文本：行用换行符分隔，列用制表符分隔。

    Args:
        table: PPT 表格对象

    Returns:
        表格文本内容
    """
    rows: List[str] = []
    for row in table.rows:
        values: List[str] = []
        for cell in row.cells:
            cell_text = " ".join(cell.text.split())
            values.append(cell_text)
        if any(values):
            rows.append("\t".join(values))
    return "\n".join(rows)


def _parse_pptx(file_path: Path, output_dir: Optional[Path]) -> List[PageContent]:
    """
    解析 .pptx 文件：每张幻灯片视为一页，提取文本、表格和图片。

    Args:
        file_path: .pptx 文件路径
        output_dir: 图片输出目录（可选）

    Returns:
        按幻灯片组织的页面内容列表
    """
    presentation = Presentation(str(file_path))
    pages: List[PageContent] = []
    image_output_dir = ensure_image_output_dir(output_dir)
    image_index = 1
    # 图片类型枚举
    picture_types = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }

    for slide_index, slide in enumerate(presentation.slides, start=1):
        # 每一张幻灯片视作一页
        texts: List[str] = []
        images: List[ImageInfo] = []
        for shape in slide.shapes:
            # 提取文本框内容
            if shape.has_text_frame:
                text = "\n".join(
                    paragraph.text for paragraph in shape.text_frame.paragraphs
                )
                if text.strip():
                    texts.append(text.strip())
            # 提取表格内容
            if shape.has_table:
                table_text = _extract_table_text(shape.table)
                if table_text.strip():
                    texts.append(table_text.strip())
            # 提取图片
            if shape.shape_type in picture_types:
                try:
                    image = shape.image
                    filename = getattr(image, "filename", None)
                    content_type = getattr(image, "content_type", None)
                    # 根据文件名或 MIME 类型确定扩展名
                    if filename and Path(filename).suffix:
                        extension = Path(filename).suffix
                    else:
                        mapping = {
                            "image/png": ".png",
                            "image/jpeg": ".jpg",
                            "image/jpg": ".jpg",
                            "image/gif": ".gif",
                            "image/bmp": ".bmp",
                            "image/tiff": ".tiff",
                            "image/webp": ".webp",
                        }
                        extension = mapping.get(content_type or "", ".img")

                    # 保存图片到输出目录
                    image_path = save_image_bytes(
                        output_dir=image_output_dir,
                        file_stem=file_path.stem,
                        page_number=slide_index,
                        index=image_index,
                        data=image.blob,
                        extension=extension,
                    )
                    images.append(
                        ImageInfo(
                            index=image_index,
                            name=filename,
                            width=getattr(image, "width", None),
                            height=getattr(image, "height", None),
                            mime_type=content_type,
                            path=str(image_path),
                        )
                    )
                    image_index += 1
                except Exception:
                    continue
        pages.append(
            PageContent(
                page_number=slide_index, text="\n".join(texts).strip(), images=images
            )
        )
    return pages


def _parse_with_unstructured(file_path: Path) -> List[PageContent]:
    """
    使用 unstructured 库降级解析旧格式 .ppt 文件。

    Args:
        file_path: .ppt 文件路径

    Returns:
        按页组织的内容列表

    Raises:
        ServiceException: 缺少依赖或解析失败
    """
    try:
        from unstructured.partition.auto import partition
    except Exception as exc:
        raise ServiceException("缺少 unstructured 依赖，无法解析该格式") from exc

    try:
        elements = partition(filename=str(file_path))
    except Exception as exc:
        raise ServiceException(f"解析文件失败: {exc}") from exc

    pages_map: Dict[int, PageContent] = {}
    for element in elements:
        # unstructured 的 page_number 可能缺失，默认归为第 1 页
        page_number = getattr(element.metadata, "page_number", None) or 1
        text = element.text or ""
        page = pages_map.setdefault(
            page_number, PageContent(page_number=page_number, text="")
        )
        if text:
            if page.text:
                page.text += "\n"
            page.text += text

    # 按页码排序
    pages = [pages_map[key] for key in sorted(pages_map)]
    return pages


class PptxLoader(FileLoader):
    """PPT 解析器，支持 pptx（优先）及 ppt（降级解析）。"""

    def parse(
            self, file_path: Path, output_dir: Optional[Path] = None
    ) -> List[PageContent]:
        """
        解析 PPT 文件。

        Args:
            file_path: PPT 文件路径
            output_dir: 图片输出目录（可选）

        Returns:
            按幻灯片组织的页面内容列表

        Raises:
            ServiceException: 不支持的文件格式
        """
        suffix = file_path.suffix.lower()
        if suffix == ".pptx":
            return _parse_pptx(file_path, output_dir)
        if suffix == ".ppt":
            return _parse_with_unstructured(file_path)
        raise ServiceException(f"不支持的 PPT 格式: {suffix}")
